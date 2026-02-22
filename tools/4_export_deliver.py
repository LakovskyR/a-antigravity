"""
Export & Deliver â€” Module 4
Exports results in selected format: JSON, WebApp, Notebook, or PowerPoint.
Uploads outputs to SharePoint.
"""

import json
import os
import subprocess
import sys
from datetime import datetime
from pathlib import Path


def load_config():
    """Load run configuration."""
    config_path = Path(__file__).parent.parent / "tmp" / "run_config.json"
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)


def save_config(config):
    """Save updated configuration."""
    config_path = Path(__file__).parent.parent / "tmp" / "run_config.json"
    with open(config_path, "w", encoding="utf-8") as f:
        json.dump(config, f, indent=2, ensure_ascii=False)


def load_final_payload():
    """Load final delivery payload."""
    payload_path = Path(__file__).parent.parent / "tmp" / "final_delivery_payload.json"
    with open(payload_path, "r", encoding="utf-8") as f:
        return json.load(f)


def export_json(payload: dict, config: dict) -> Path:
    """
    Export as JSON to SharePoint/Analytics folder.
    Also saves locally to tmp/outputs/.
    """
    print("   ðŸ“„ Exporting JSON...")
    
    meta = config["project_metadata"]
    project_name = meta["project_name"]
    wave = meta["wave"]
    
    # Local output
    output_dir = Path(__file__).parent.parent / "tmp" / "outputs"
    output_dir.mkdir(parents=True, exist_ok=True)
    
    filename = f"{project_name}_{wave}_analysis.json"
    output_path = output_dir / filename
    
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, indent=2, ensure_ascii=False)
    
    print(f"      ðŸ’¾ Saved: {output_path}")
    
    # TODO: Upload to SharePoint
    sharepoint_path = f"/Analytics/{project_name}/{wave}/{filename}"
    print(f"      â¬†ï¸  SharePoint target: {sharepoint_path}")
    
    return output_path


def export_webapp(payload: dict, config: dict) -> Path:
    """
    Launch Streamlit web application.
    """
    print("   ðŸŒ Launching Streamlit WebApp...")
    
    # Copy payload to location accessible by app.py
    output_dir = Path(__file__).parent.parent / "tmp" / "outputs"
    output_dir.mkdir(parents=True, exist_ok=True)
    
    payload_path = output_dir / "current_payload.json"
    with open(payload_path, "w", encoding="utf-8") as f:
        json.dump(payload, f, indent=2, ensure_ascii=False)
    
    print(f"      ðŸ’¾ Payload saved: {payload_path}")
    
    # Check if app.py exists
    app_path = Path(__file__).parent.parent / "app.py"
    if app_path.exists():
        print(f"      ðŸš€ Starting Streamlit...")
        print(f"      Run manually: streamlit run app.py")
        # Note: We don't auto-launch to avoid blocking the pipeline
    else:
        print(f"      âš ï¸  app.py not found. Phase 4 will build it.")
    
    return payload_path


def export_notebook(payload: dict, config: dict) -> Path:
    """
    Generate Jupyter Notebook from payload using nbformat.
    """
    print("   ðŸ““ Generating Jupyter Notebook...")
    
    try:
        import nbformat as nbf
    except ImportError:
        print("      âš ï¸  nbformat not installed. Installing...")
        subprocess.run([sys.executable, "-m", "pip", "install", "nbformat"], 
                      capture_output=True)
        import nbformat as nbf
    
    meta = config["project_metadata"]
    project_name = meta["project_name"]
    wave = meta["wave"]
    
    # Create notebook
    nb = nbf.v4.new_notebook()
    
    # Title cell
    title = f"# {project_name} â€” {wave} Analysis Report\n\n"
    title += f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n"
    title += f"Modules: {', '.join(meta['selected_modules'])}\n\n"
    title += f"LLM Provider: {meta['llm_provider']}"
    
    nb.cells.append(nbf.v4.new_markdown_cell(title))
    
    # Add cells for each module
    for item in payload.get("analytics_payload", []):
        module_name = item.get("analysis_type", "Unknown")
        
        # Module header
        nb.cells.append(nbf.v4.new_markdown_cell(f"## {module_name.title()} Analysis"))
        
        # AI Insights
        insights = item.get("ai_insights", {})
        academic = insights.get("academic_commentary", "N/A")
        business = insights.get("business_translation", "N/A")
        
        insights_md = f"""### Academic Perspective

{academic}

### Business Perspective

{business}
"""
        nb.cells.append(nbf.v4.new_markdown_cell(insights_md))
        
        # Chart data visualization code
        chart_data = item.get("chart_data", {})
        if chart_data.get("labels"):
            code = f"""import matplotlib.pyplot as plt

# {module_name} visualization
labels = {chart_data.get('labels', [])}
values = {chart_data.get('values', [])}

plt.figure(figsize=(10, 6))
plt.bar(labels, values)
plt.title('{module_name.title()} Results')
plt.xlabel('Categories')
plt.ylabel('Count')
plt.xticks(rotation=45)
plt.tight_layout()
plt.show()
"""
            nb.cells.append(nbf.v4.new_code_cell(code))
    
    # Summary cell
    summary = payload.get("summary", {})
    recommendations = summary.get("key_recommendations", [])
    
    summary_md = "## Summary & Recommendations\n\n"
    summary_md += f"**Confidence Level:** {summary.get('confidence_level', 'N/A')}\n\n"
    summary_md += "**Key Recommendations:**\n\n"
    for rec in recommendations:
        summary_md += f"- {rec}\n"
    
    nb.cells.append(nbf.v4.new_markdown_cell(summary_md))
    
    # Save notebook
    output_dir = Path(__file__).parent.parent / "tmp" / "outputs"
    output_dir.mkdir(parents=True, exist_ok=True)
    
    filename = f"{project_name}_{wave}_analysis.ipynb"
    output_path = output_dir / filename
    
    with open(output_path, "w", encoding="utf-8") as f:
        nbf.write(nb, f)
    
    print(f"      ðŸ’¾ Saved: {output_path}")
    
    # SharePoint target
    sharepoint_path = f"/Analytics/{project_name}/{wave}/{filename}"
    print(f"      â¬†ï¸  SharePoint target: {sharepoint_path}")
    
    return output_path


def export_powerpoint(payload: dict, config: dict) -> Path:
    """
    Generate PowerPoint from payload using python-pptx.
    """
    print("   ðŸ“Š Generating PowerPoint...")
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("      âš ï¸  python-pptx not installed. Installing...")
        subprocess.run([sys.executable, "-m", "pip", "install", "python-pptx"], 
                      capture_output=True)
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
    
    meta = config["project_metadata"]
    project_name = meta["project_name"]
    wave = meta["wave"]
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"{project_name}\n{wave} Analysis"
    title_slide.placeholders[1].text = f"Generated: {datetime.now().strftime('%Y-%m-%d')}\nModules: {', '.join(meta['selected_modules'])}"
    
    # Add slides for each module
    for item in payload.get("analytics_payload", []):
        module_name = item.get("analysis_type", "Unknown").title()
        insights = item.get("ai_insights", {})
        
        # Module slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = module_name
        
        # Academic commentary
        tf = slide.placeholders[1].text_frame
        tf.text = "Academic Perspective:"
        p = tf.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        
        academic = insights.get("academic_commentary", "N/A")
        tf.add_paragraph().text = academic[:500] + "..." if len(academic) > 500 else academic
        
        # Business commentary slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = f"{module_name} â€” Business View"
        
        tf2 = slide2.placeholders[1].text_frame
        tf2.text = "Business Perspective:"
        p2 = tf2.paragraphs[0]
        p2.font.bold = True
        p2.font.size = Pt(14)
        
        business = insights.get("business_translation", "N/A")
        tf2.add_paragraph().text = business[:500] + "..." if len(business) > 500 else business
    
    # Summary slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_slide.shapes.title.text = "Summary & Recommendations"
    
    summary = payload.get("summary", {})
    tf = summary_slide.placeholders[1].text_frame
    tf.text = f"Confidence Level: {summary.get('confidence_level', 'N/A')}"
    
    for rec in summary.get("key_recommendations", []):
        tf.add_paragraph().text = f"â€¢ {rec}",
    
    # Save presentation
    output_dir = Path(__file__).parent.parent / "tmp" / "outputs"
    output_dir.mkdir(parents=True, exist_ok=True)
    
    filename = f"{project_name}_{wave}_analysis.pptx"
    output_path = output_dir / filename
    
    prs.save(output_path)
    
    print(f"      ðŸ’¾ Saved: {output_path}")
    
    # Canva integration note
    print(f"      ðŸ’¡ Canva API integration: Template population available if API key configured")
    
    # SharePoint target
    sharepoint_path = f"/Analytics/{project_name}/{wave}/{filename}"
    print(f"      â¬†ï¸  SharePoint target: {sharepoint_path}")
    
    return output_path


def upload_to_sharepoint(file_path: Path, sharepoint_path: str) -> bool:
    """
    Upload file to SharePoint.
    TODO: Implement Microsoft Graph API upload.
    """
    print(f"      â¬†ï¸  Uploading to SharePoint: {sharepoint_path}")
    print(f"      âš ï¸  SharePoint upload not yet implemented")
    print(f"      ðŸ’¡ File available locally: {file_path}")
    return True


def main():
    """Main export and delivery workflow."""
    print("=" * 70)
    print("ðŸ“¤ MODULE 4: EXPORT & DELIVER")
    print("=" * 70)
    
    config = load_config()
    payload = load_final_payload()
    
    delivery_format = config["project_metadata"]["delivery_format"]
    print(f"\nðŸŽ¯ Delivery format: {delivery_format}")
    
    # Update status
    config["run_info"]["status"] = "exporting"
    save_config(config)
    
    # Route to appropriate export function
    export_functions = {
        "json": export_json,
        "webapp": export_webapp,
        "notebook": export_notebook,
        "powerpoint": export_powerpoint
    }
    
    if delivery_format in export_functions:
        output_path = export_functions[delivery_format](payload, config)
    else:
        print(f"   âŒ Unknown delivery format: {delivery_format}")
        sys.exit(1)
    
    # Update config
    config["run_info"]["status"] = "exported"
    config["run_info"]["output_file"] = str(output_path)
    save_config(config)
    
    print("\n" + "=" * 70)
    print("âœ… EXPORT & DELIVER COMPLETE")
    print("=" * 70)


if __name__ == "__main__":
    main()

